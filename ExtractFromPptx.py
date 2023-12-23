import tkinter as tk
from tkinter import filedialog, simpledialog
from pptx import Presentation
import os

def extract_text_from_pptx(file_path, slide_numbers):
    ppt = Presentation(file_path)
    extracted_text = ""

    for i in slide_numbers:
        slide = ppt.slides[i - 1]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"

    return extracted_text

def parse_slide_numbers(input_str):
    ranges = (number_range.split("-") for number_range in input_str.split(","))
    return [i for r in ranges for i in range(int(r[0]), int(r[-1]) + 1)]

def select_file():
    current_directory = os.getcwd()
    file_path = filedialog.askopenfilename(initialdir=current_directory,
                                           filetypes=[("PowerPoint files", "*.pptx")])
    if file_path:
        slide_numbers_input = simpledialog.askstring("Input", "Enter slide numbers or ranges (e.g., 1-3,5,7-9)")
        slide_numbers = parse_slide_numbers(slide_numbers_input)
        extracted_text = extract_text_from_pptx(file_path, slide_numbers)
        text_box.delete('1.0', tk.END)
        text_box.insert(tk.END, extracted_text)

root = tk.Tk()
root.title("PPTX Text Extractor")

open_button = tk.Button(root, text="Open PPTX File", command=select_file)
open_button.pack()

text_box = tk.Text(root, height=15, width=50)
text_box.pack()

root.mainloop()
